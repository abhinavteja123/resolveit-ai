"""Generate ResolveIT_AI_Presentation.pptx — visually matched to the
website's Tailwind theme (dark brown background, amber primary, glass-card
surfaces).

Palette (lifted from frontend/tailwind.config.js + frontend/src/index.css):
  body              #0e0b09  (dark-950)
  card              #1a1510  (dark-900)
  card-stroke       #2a221c  (dark-800)
  text-primary      #e8dfd6  (dark-200)
  text-muted        #9c8b80  (dark-500)
  text-faint        #6d5f57  (dark-600)
  accent-amber      #f59e0b  (primary-500)
  accent-amber-dark #d97706  (primary-600)
  accent-amber-soft #fbbf24  (primary-400)
  emerald-500       #10b981
  amber-400         #fbbf24
  red-500           #ef4444
  sky-500           #0ea5e9
  font UI           Plus Jakarta Sans (falls back to Calibri)
  font code         JetBrains Mono
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# ── Palette ────────────────────────────────────────────────────────
BG        = RGBColor(0x0e, 0x0b, 0x09)
CARD      = RGBColor(0x1a, 0x15, 0x10)
CARD_LITE = RGBColor(0x2a, 0x22, 0x1c)
STROKE    = RGBColor(0x3d, 0x33, 0x2c)
TEXT      = RGBColor(0xe8, 0xdf, 0xd6)
TEXT_MUTE = RGBColor(0x9c, 0x8b, 0x80)
TEXT_FAINT= RGBColor(0x6d, 0x5f, 0x57)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
AMBER_D   = RGBColor(0xd9, 0x77, 0x06)
AMBER_L   = RGBColor(0xfb, 0xbf, 0x24)
EMERALD   = RGBColor(0x10, 0xb9, 0x81)
RED       = RGBColor(0xef, 0x44, 0x44)
SKY       = RGBColor(0x0e, 0xa5, 0xe9)
WHITE     = RGBColor(0xff, 0xff, 0xff)

UI_FONT   = "Plus Jakarta Sans"   # falls back gracefully on systems w/o it
MONO_FONT = "JetBrains Mono"


# ── Helpers ────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, italic=False, color=TEXT,
             align=PP_ALIGN.LEFT, font=UI_FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0;  tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=13, color=TEXT, accent=AMBER, gap=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = gap
        rb = p.add_run()
        rb.text = "›  "
        rb.font.name = UI_FONT
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = accent
        r = p.add_run()
        r.text = it
        r.font.name = UI_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def card(slide, left, top, width, height, *, fill=CARD, stroke=CARD_LITE, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(0.75)
    return shp


def chip(slide, left, top, width, height, text, *,
         fill=None, stroke=None, color=TEXT, size=11, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    if fill is None:
        shp.fill.solid(); shp.fill.fore_color.rgb = CARD
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke if stroke is not None else CARD_LITE
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04);  tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = UI_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shp


def amber_chip(slide, left, top, width, height, text, *, size=11, bold=True):
    """Glow-style amber chip used as the accent throughout the deck."""
    shp = chip(slide, left, top, width, height, text,
               fill=CARD, stroke=AMBER, color=AMBER_L,
               size=size, bold=bold)
    return shp


def header(slide, slide_num, total, kicker, title):
    # Top kicker (orange small caps)
    add_text(slide, Inches(0.5), Inches(0.32), Inches(8), Inches(0.3),
             kicker.upper(), size=10, bold=True, color=AMBER)
    # Big title
    add_text(slide, Inches(0.5), Inches(0.55), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=TEXT)
    # Underline divider
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.18),
                                      Inches(2.5), Inches(1.18))
    line.line.color.rgb = AMBER
    line.line.width = Pt(2.0)
    # Slide number top-right
    add_text(slide, Inches(11.8), Inches(0.45), Inches(1.3), Inches(0.3),
             f"{slide_num} / {total}",
             size=11, bold=True, color=TEXT_MUTE, align=PP_ALIGN.RIGHT)


def footer(slide):
    # Thin divider above footer
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.05),
                                      Inches(12.83), Inches(7.05))
    line.line.color.rgb = STROKE
    line.line.width = Pt(0.5)
    add_text(slide, Inches(0.5), Inches(7.13), Inches(6), Inches(0.3),
             "ResolveIT AI · Smart Runbook Resolution Assistant",
             size=9, italic=True, color=TEXT_FAINT)


# Add an "ambient orb" — large blurred orange shape — for visual depth.
# python-pptx doesn't support blur, so we use a low-opacity circle.
def ambient(slide, left, top, diameter, *, color=AMBER, opacity_pct=8):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    # Apply alpha via XML
    sp = shp.fill.fore_color._xFill
    # Wrap solid-fill colour with alpha
    srgb = sp.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', str(int(opacity_pct * 1000)))
    return shp


def build_table(slide, left, top, width, height, headers, rows,
                *, size=11, header_size=11, header_fill=AMBER, alt=True):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                       left, top, width, height)
    tbl = tbl_shape.table
    # Header row
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        c.fill.solid()
        c.fill.fore_color.rgb = header_fill
        c.text_frame.text = ""
        p = c.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.name = UI_FONT
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = BG
    # Body rows
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if (ri % 2 == 1 or not alt) else BG
            c.text_frame.text = ""
            p = c.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.name = UI_FONT
            r.font.size = Pt(size)
            r.font.color.rgb = TEXT
    # Style cell borders
    for row in tbl.rows:
        for c in row.cells:
            tcPr = c._tc.get_or_add_tcPr()
            for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
                ln = etree.SubElement(tcPr, qn(f'a:{edge}'),
                                      attrib={'w': '6350', 'cap': 'flat',
                                              'cmpd': 'sng', 'algn': 'ctr'})
                fill = etree.SubElement(ln, qn('a:solidFill'))
                clr = etree.SubElement(fill, qn('a:srgbClr'),
                                       attrib={'val': '2a221c'})
    return tbl


def feature_card(slide, left, top, width, height, title, body, *,
                 accent=AMBER):
    shp = card(slide, left, top, width, height)
    # Left accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, Inches(0.06), height)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12);  tf.margin_bottom = Inches(0.12)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = UI_FONT; r1.font.size = Pt(13); r1.font.bold = True
    r1.font.color.rgb = TEXT
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.15
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = UI_FONT; r2.font.size = Pt(10)
    r2.font.color.rgb = TEXT_MUTE


def connector(slide, x1, y1, x2, y2, *, color=STROKE, width_pt=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


# ────────────────────────────────────────────────────────────────────
# Build presentation
# ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
TOTAL = 10


# ── Slide 1: Title ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)
ambient(s, Inches(-2),  Inches(-1.5), Inches(8), color=AMBER, opacity_pct=8)
ambient(s, Inches(8),   Inches(4.5),  Inches(7), color=AMBER_D, opacity_pct=6)
# Brand kicker
add_text(s, Inches(0.5), Inches(2.05), Inches(12.3), Inches(0.4),
         "INTERNAL IT SUPPORT  ·  RAG  ·  STREAMING", size=12, bold=True,
         color=AMBER, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
         "ResolveIT AI", size=72, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
# Tagline
add_text(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.55),
         "Smart Runbook Resolution Assistant", size=22, italic=True,
         color=TEXT_MUTE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.45),
         "Cited answers · Six modes · Hybrid retrieval · Real-time streaming",
         size=15, color=TEXT, align=PP_ALIGN.CENTER)

# Tech chips
chips_text = ["FastAPI", "React + Vite", "FAISS + BM25", "BGE Reranker",
              "Gemini 2.5 Flash", "Supabase", "Firebase Auth"]
chip_w = Inches(1.55); gap = Inches(0.18)
total_w = chip_w * len(chips_text) + gap * (len(chips_text) - 1)
start = (prs.slide_width - total_w) // 2
for i, t in enumerate(chips_text):
    chip(s, start + i * (chip_w + gap), Inches(5.35), chip_w, Inches(0.42),
         t, fill=CARD, stroke=AMBER, color=AMBER_L, size=10)

# Bottom strap
add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
         "Project Documentation Deck  ·  10 slides", size=11,
         italic=True, color=TEXT_FAINT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
         "ResolveIT AI · Slide 1 / 10", size=9, italic=True,
         color=TEXT_FAINT, align=PP_ALIGN.CENTER)


# ── Slide 2: Problem & Motivation ──────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
ambient(s, Inches(-2), Inches(-1.5), Inches(6), color=AMBER, opacity_pct=6)
header(s, 2, TOTAL, "The Problem", "Why ResolveIT AI exists")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
         "Internal IT teams keep runbooks scattered across PDFs, wikis, and shared drives — finding the right fix during an incident is slow, error-prone, and ungrounded.",
         size=13, italic=True, color=TEXT_MUTE)

# Two cards
card(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(4.7))
add_text(s, Inches(0.7), Inches(2.15), Inches(5.8), Inches(0.5),
         "Pain Points", size=18, bold=True, color=AMBER)
add_bullets(s, Inches(0.7), Inches(2.7), Inches(5.8), Inches(4.0), [
    "Engineers waste minutes-to-hours per incident searching for fixes.",
    "Keyword search returns documents, not answers.",
    "Generic chatbots hallucinate — unsafe for production actions.",
    "Senior SREs and new-hires need very different answer styles.",
    "Knowledge gaps stay invisible until they cause an outage.",
], size=13, color=TEXT, accent=RED)

card(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.7))
add_text(s, Inches(7.05), Inches(2.15), Inches(5.7), Inches(0.5),
         "Goals", size=18, bold=True, color=AMBER)
add_bullets(s, Inches(7.05), Inches(2.7), Inches(5.7), Inches(4.0), [
    "Citation-grounded resolutions tied to internal runbooks.",
    "Mode-aware answers that adapt to user expertise.",
    "Real-time streaming with verifiable inline sources.",
    "Closed-loop feedback so admins detect knowledge gaps.",
    "Production-ready: auth, rate limits, cache, idempotent migrations.",
], size=13, color=TEXT, accent=EMERALD)
footer(s)


# ── Slide 3: Solution & Key Features ───────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
ambient(s, Inches(10), Inches(-2), Inches(6), color=AMBER, opacity_pct=6)
header(s, 3, TOTAL, "Solution Overview", "Key Capabilities at a Glance")

cards_8 = [
    ("Hybrid Retrieval",
     "FAISS dense + BM25 lexical, fused via Reciprocal Rank Fusion (k=60)."),
    ("HyDE Expansion",
     "Gemini writes a hypothetical passage to bridge query/document vocabulary."),
    ("Cross-Encoder Re-rank",
     "BGE-reranker-base scores (query, chunk) pairs; sigmoid → confidence ∈ [0,1]."),
    ("Strict Gate Check",
     "If no excerpt directly addresses the query, refuses with an escalation sentinel."),
    ("Six Response Modes",
     "Fast / Standard / Deep / ELI5 / Expert / Dry-run."),
    ("Streaming SSE Answers",
     "Tokens stream in real time; events: mode → sources → token → done."),
    ("Regenerate · Follow-ups · Export",
     "Re-run in any mode; AI-suggested next questions; download as Markdown."),
    ("Threading · Bookmarks · Sharing",
     "thread_id continuation; Playbook saves; /answer/:id permalinks."),
]
cw = Inches(2.95); ch = Inches(1.35); gx = Inches(0.18); gy = Inches(0.18)
ox = Inches(0.5);  oy = Inches(1.55)
for i, (t, b) in enumerate(cards_8):
    col = i % 4; row = i // 4
    feature_card(s, ox + col*(cw+gx), oy + row*(ch+gy), cw, ch, t, b)

# Stat strip
strip_y = Inches(4.85)
stats = [("3", "ML Models"),
         ("6", "Response Modes"),
         ("4", "DB Tables"),
         ("20", "API Routes"),
         ("13", "Pipeline Stages"),
         ("60", "RRF Constant")]
sw = Inches(2.0); sg = Inches(0.10)
total_sw = sw*len(stats) + sg*(len(stats)-1)
sx = (prs.slide_width - total_sw) // 2
for i, (n, l) in enumerate(stats):
    bx = sx + i*(sw+sg)
    card(s, bx, strip_y, sw, Inches(1.6))
    add_text(s, bx, strip_y + Inches(0.18), sw, Inches(0.7),
             n, size=36, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_text(s, bx, strip_y + Inches(0.95), sw, Inches(0.4),
             l, size=11, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
footer(s)


# ── Slide 4: Tech Stack & Models ───────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
header(s, 4, TOTAL, "Stack", "Tech Stack & Models Used")

card(s, Inches(0.45), Inches(1.4), Inches(6.15), Inches(5.4))
add_text(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(0.4),
         "Backend", size=15, bold=True, color=AMBER)
build_table(s, Inches(0.6), Inches(2.0), Inches(5.85), Inches(4.7),
            ["Layer", "Technology"], [
                ["Web framework", "FastAPI 0.115 (async + SSE)"],
                ["LLM", "Gemini 2.5 Flash (+ fallback list)"],
                ["Embedding", "BAAI/bge-small-en-v1.5 (384-d)"],
                ["Re-ranker", "BAAI/bge-reranker-base (cross-encoder)"],
                ["Dense store", "FAISS IndexFlatIP (cosine via L2-norm)"],
                ["Lexical store", "BM25Okapi (rank-bm25)"],
                ["Database", "Supabase (PostgreSQL + REST)"],
                ["Auth verify", "firebase-admin SDK (JWT)"],
                ["Parsing", "PyMuPDF · python-docx"],
                ["Rate · Cache", "SlowAPI · cachetools TTLCache"],
            ], size=10, header_size=10)

card(s, Inches(6.75), Inches(1.4), Inches(6.1), Inches(5.4))
add_text(s, Inches(6.9), Inches(1.55), Inches(5.85), Inches(0.4),
         "Frontend & Infra", size=15, bold=True, color=AMBER)
build_table(s, Inches(6.9), Inches(2.0), Inches(5.85), Inches(4.7),
            ["Layer", "Technology"], [
                ["UI framework", "React 18.3"],
                ["Bundler", "Vite 5"],
                ["Routing", "react-router-dom 6"],
                ["HTTP / SSE", "Axios + native fetch"],
                ["Auth client", "Firebase JS SDK 10"],
                ["Styling", "Tailwind CSS 3.4"],
                ["Animations", "Framer Motion"],
                ["Markdown render", "react-markdown"],
                ["Containers", "Docker + Docker Compose"],
                ["Auth provider", "Firebase Authentication"],
            ], size=10, header_size=10)
footer(s)


# ── Slide 5: System Architecture ───────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
ambient(s, Inches(-2), Inches(2), Inches(6), color=AMBER, opacity_pct=5)
header(s, 5, TOTAL, "Architecture", "System Architecture — Four Layers")

layer_left = Inches(0.5); layer_w = Inches(12.35); layer_h = Inches(1.05)
y_client = Inches(1.45); y_api = Inches(2.65)
y_svc    = Inches(3.85); y_data = Inches(5.05)


def layer(top, label):
    card(s, layer_left, top, layer_w, layer_h, fill=CARD, stroke=CARD_LITE)
    # Left tab
    tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             layer_left, top, Inches(0.08), layer_h)
    tab.fill.solid(); tab.fill.fore_color.rgb = AMBER
    tab.line.fill.background()
    add_text(s, layer_left + Inches(0.18), top + Inches(0.05),
             Inches(2.5), Inches(0.4),
             label, size=12, bold=True, color=AMBER)


layer(y_client, "CLIENT")
layer(y_api,    "API")
layer(y_svc,    "SERVICE")
layer(y_data,   "DATA · EXTERNAL")


def inner(x, y, w, h, text):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = BG
    shp.line.color.rgb = STROKE; shp.line.width = Pt(0.5)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = UI_FONT; r.font.size = Pt(9.5); r.font.color.rgb = TEXT
    return shp


bw = Inches(2.2); bh = Inches(0.7); ix = Inches(3.05); gap = Inches(0.18)


def row_boxes(y_layer, items):
    boxes = []
    for i, t in enumerate(items):
        boxes.append(inner(ix + i*(bw+gap), y_layer + Inches(0.22), bw, bh, t))
    return boxes


cli = row_boxes(y_client, ["Browser\nReact 18 SPA",
                           "AuthContext\n(Firebase token)",
                           "useQuery hook\n(fetch SSE)",
                           "Pages\nDashboard · Admin …"])
api = row_boxes(y_api, ["/query · /query/stream\n(SSE)",
                        "/auth · /export\n/feedback · /history",
                        "/admin/* · /runbooks\n/bookmarks",
                        "Middleware\nSlowAPI · CORS · JWT · Cache"])
svc = row_boxes(y_svc, ["RAG Pipeline\nrag/pipeline.py",
                        "Modes Pack\nrag/modes.py",
                        "Retrieval\nFAISS · BM25 · Reranker",
                        "Ingestion\nparser · chunker · indexer"])
dat = row_boxes(y_data, ["Supabase\nPostgreSQL",
                         "FAISS index.bin\n+ metadata.json",
                         "Firebase Auth\nOAuth + JWT",
                         "Gemini API\n+ Hugging Face"])

# Vertical amber connectors between layers (centered)
mid = ix + 2*(bw+gap) + bw // 2
for top, bot in [(cli, api), (api, svc), (svc, dat)]:
    src = top[2]; dst = bot[0]
    connector(s,
              src.left + src.width // 2, src.top + src.height,
              dst.left + dst.width // 2, dst.top,
              color=AMBER, width_pt=1.5)

add_text(s, Inches(0.5), Inches(6.35), Inches(12.35), Inches(0.4),
         "Client · API · Service · Data — clean separation; only the API layer holds JWT verify and rate-limit.",
         size=11, italic=True, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
footer(s)


# ── Slide 6: RAG Pipeline (full diagram) ───────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
ambient(s, Inches(10), Inches(4), Inches(6), color=AMBER_D, opacity_pct=6)
header(s, 6, TOTAL, "Pipeline", "Full RAG Pipeline — Query → Cited Answer")

bw = Inches(1.55); bh = Inches(0.78); gap = Inches(0.16)
sx = Inches(0.5)


def step(x, y, lbl, *, accent=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh)
    shp.fill.solid()
    shp.fill.fore_color.rgb = CARD
    shp.line.color.rgb = AMBER if accent else CARD_LITE
    shp.line.width = Pt(1.5 if accent else 0.75)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lbl
    r.font.name = UI_FONT; r.font.size = Pt(9.5); r.font.bold = True
    r.font.color.rgb = AMBER_L if accent else TEXT
    return shp


def arrow(a, b, *, color=AMBER):
    line = s.shapes.add_connector(1,
                                  a.left + a.width, a.top + a.height // 2,
                                  b.left, b.top + b.height // 2)
    line.line.color.rgb = color; line.line.width = Pt(1.25)
    # add arrow head
    ln = line.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'sm'); tail.set('len', 'sm')
    return line


row1_y = Inches(1.65); row2_y = Inches(3.55)
row1 = ["Query", "Auth +\nMode Lookup", "Cache Check\n(mode-keyed)",
        "HyDE Expand\n(Gemini)", "Embed\nBGE-small",
        "Hybrid\nFAISS+BM25", "RRF Fuse"]
boxes_1 = [step(sx + i*(bw+gap), row1_y, t,
                accent=(t in ("HyDE Expand\n(Gemini)",
                              "Hybrid\nFAISS+BM25")))
           for i, t in enumerate(row1)]
for i in range(len(boxes_1) - 1):
    arrow(boxes_1[i], boxes_1[i+1])

row2 = ["Scope Filter\nadmin · mine · both",
        "Re-rank\nCross-Encoder", "Gate\n(score ≥ 0.25)",
        "Gemini\nstream tokens", "Follow-ups\n(Gemini × 3)",
        "Log + Cache", "SSE 'done'"]
boxes_2 = [step(sx + i*(bw+gap), row2_y, t,
                accent=(t in ("Re-rank\nCross-Encoder",
                              "Gemini\nstream tokens")))
           for i, t in enumerate(row2)]
for i in range(len(boxes_2) - 1):
    arrow(boxes_2[i], boxes_2[i+1])

# Down-then-back link between rows
last1 = boxes_1[-1]; first2 = boxes_2[0]
mid_x_a = last1.left + last1.width // 2
mid_x_b = first2.left + first2.width // 2
mid_y = (row1_y + Inches(0.78) + row2_y) // 2
# down from last of row1
l1 = s.shapes.add_connector(1, mid_x_a, row1_y + bh, mid_x_a, mid_y)
l1.line.color.rgb = AMBER; l1.line.width = Pt(1.25)
# left across
l2 = s.shapes.add_connector(1, mid_x_a, mid_y, mid_x_b, mid_y)
l2.line.color.rgb = AMBER; l2.line.width = Pt(1.25)
# up into first of row2
l3 = s.shapes.add_connector(1, mid_x_b, mid_y, mid_x_b, row2_y)
l3.line.color.rgb = AMBER; l3.line.width = Pt(1.25)
# arrow head on last segment
for ln_shp in (l3,):
    ln = ln_shp.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'sm'); tail.set('len', 'sm')

# Notes block
note_y = Inches(5.0)
card(s, Inches(0.5), note_y, Inches(12.35), Inches(1.85))
add_text(s, Inches(0.7), note_y + Inches(0.12),
         Inches(11.9), Inches(0.4),
         "Mode controls these knobs",
         size=12, bold=True, color=AMBER)

knobs = [("use_hyde", "Skip in Fast"),
         ("top_k", "6 / 12 / 20"),
         ("top_n", "3 / 5 / 8"),
         ("temperature", "0.0–0.3"),
         ("system_prompt", "per-mode"),
         ("critique_retry", "Deep only")]
kx = Inches(0.7); ky = note_y + Inches(0.55)
for i, (k, v) in enumerate(knobs):
    bx = kx + i * Inches(2.0)
    add_text(s, bx, ky, Inches(2.0), Inches(0.3),
             k, size=11, bold=True, color=AMBER_L, font=MONO_FONT)
    add_text(s, bx, ky + Inches(0.3), Inches(2.0), Inches(0.3),
             v, size=10, color=TEXT_MUTE)

add_text(s, Inches(0.7), note_y + Inches(1.3), Inches(11.9), Inches(0.3),
         "SSE events:  mode  →  sources  →  token (× N)  →  done {query_log_id, top_confidence, mode, follow_ups}",
         size=10, italic=True, color=TEXT_MUTE, font=MONO_FONT)
footer(s)


# ── Slide 7: Six Response Modes ────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
header(s, 7, TOTAL, "Modes", "Six Response Modes — One Pipeline, Six Audiences")

add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
         "Each mode bundles a retrieval-config + a system prompt. Cache key includes mode → independent caching per mode.",
         size=12, italic=True, color=TEXT_MUTE)

# Mode cards in a 3 × 2 grid
modes = [
    ("Fast",     "OFF · top_k 6 · top_n 3 · 0.0",   "2 AM incident — terse, ~1 s",                  AMBER),
    ("Standard", "ON · top_k 12 · top_n 5 · 0.2",   "Default — Summary + Steps + Prevention",       EMERALD),
    ("Deep",     "ON · top_k 20 · top_n 8 · 0.2",   "Tricky bug — Root Cause + Verification",       SKY),
    ("ELI5",     "ON · top_k 12 · top_n 5 · 0.3",   "Junior tech — explains jargon + commands",     AMBER_L),
    ("Expert",   "ON · top_k 12 · top_n 5 · 0.1",   "Senior SRE — TL;DR + command-only steps",      RED),
    ("Dry-run",  "ON · top_k 12 · top_n 5 · 0.1",   "Before destructive ops — annotated + Rollback", AMBER_D),
]
cw = Inches(4.0); ch = Inches(2.1); gx = Inches(0.18); gy = Inches(0.2)
ox = Inches(0.5);  oy = Inches(1.95)
for i, (name, knobs_s, use, accent) in enumerate(modes):
    col = i % 3; row = i // 3
    x = ox + col*(cw+gx); y = oy + row*(ch+gy)
    shp = card(s, x, y, cw, ch)
    # Top accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # Title
    add_text(s, x + Inches(0.2), y + Inches(0.18),
             cw - Inches(0.4), Inches(0.5),
             name, size=22, bold=True, color=accent)
    # Knobs
    add_text(s, x + Inches(0.2), y + Inches(0.85),
             cw - Inches(0.4), Inches(0.4),
             knobs_s, size=11, color=TEXT, font=MONO_FONT, bold=True)
    # Use case
    add_text(s, x + Inches(0.2), y + Inches(1.25),
             cw - Inches(0.4), Inches(0.7),
             use, size=12, color=TEXT_MUTE)

footer(s)


# ── Slide 8: Differentiating Features ──────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
header(s, 8, TOTAL, "Beyond Vanilla RAG", "What Makes ResolveIT AI Different")

cards_12 = [
    ("Regenerate in Any Mode",
     "Re-run the same query in another mode; new log row carries regenerate_of FK."),
    ("Follow-up Question Chips",
     "Gemini proposes 3 short next questions; clicking opens a threaded continuation."),
    ("Markdown Export",
     "GET /export/{id}.md downloads a portable file with title, mode, confidence, sources."),
    ("Citation Gate-Check",
     "If no excerpt directly addresses the query, an escalation sentinel is returned."),
    ("Threaded Conversations",
     "thread_id links continuation queries; /history groups by thread; Resume re-seeds chat."),
    ("Bookmarks (Playbook)",
     "Per-user saved answers with snippet + sources, paginated."),
    ("Shareable Permalinks",
     "/answer/:id renders a read-only view of any logged answer."),
    ("Admin Knowledge Gaps",
     "Surfaces low-confidence and negative-feedback queries grouped by topic."),
    ("Runbook Health Grid",
     "Per-runbook query count, average confidence, thumbs ratio, Needs-Attention flag."),
    ("Streaming with Abort",
     "useQuery uses AbortController so a new query cancels in-flight stream cleanly."),
    ("Per-Mode TTL Cache",
     "Same query in different modes is cached independently (sha256 of mode:scope:user:query)."),
    ("Cmd+K Command Palette",
     "Keyboard-driven nav with grouped suggestions / navigation / actions."),
]
cw = Inches(2.95); ch = Inches(1.4); gx = Inches(0.2); gy = Inches(0.2)
ox = Inches(0.5);  oy = Inches(1.45)
for i, (t, b) in enumerate(cards_12):
    col = i % 4; row = i // 4
    x = ox + col*(cw+gx); y = oy + row*(ch+gy)
    feature_card(s, x, y, cw, ch, t, b)
footer(s)


# ── Slide 9: Database & API ────────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
header(s, 9, TOTAL, "Schema · API", "Database Schema & API Surface")

card(s, Inches(0.45), Inches(1.4), Inches(6.1), Inches(5.4))
add_text(s, Inches(0.6), Inches(1.55), Inches(5.85), Inches(0.4),
         "Supabase Tables", size=15, bold=True, color=AMBER)
build_table(s, Inches(0.6), Inches(2.0), Inches(5.85), Inches(4.5),
            ["Table", "Key Columns"], [
                ["runbooks", "id · filename · category · file_type · uploaded_by · content_hash · is_admin_runbook"],
                ["query_logs", "id · user_id · query_text · sources · llm_response · confidence · thread_id · mode · regenerate_of"],
                ["feedback", "id · query_log_id · user_id · rating (1/-1) · comment"],
                ["bookmarks", "id · user_id · query_log_id · query_text · answer_snippet · sources"],
            ], size=10, header_size=10)
add_text(s, Inches(0.6), Inches(6.5), Inches(5.85), Inches(0.3),
         "Self-FKs: query_logs.thread_id → query_logs.id  ·  query_logs.regenerate_of → query_logs.id",
         size=9, italic=True, color=TEXT_FAINT, font=MONO_FONT)

card(s, Inches(6.75), Inches(1.4), Inches(6.1), Inches(5.4))
add_text(s, Inches(6.9), Inches(1.55), Inches(5.85), Inches(0.4),
         "Key API Routes", size=15, bold=True, color=AMBER)
build_table(s, Inches(6.9), Inches(2.0), Inches(5.85), Inches(4.7),
            ["Method", "Path"], [
                ["POST", "/query  ·  /query/stream  (SSE)"],
                ["GET",  "/answer/{id}        — permalink"],
                ["GET",  "/export/{id}.md     — Markdown download"],
                ["POST", "/feedback           — thumbs + comment"],
                ["GET",  "/history            — threaded log"],
                ["POST · GET · DELETE", "/bookmarks"],
                ["POST · GET", "/runbooks/upload · /runbooks/my"],
                ["POST · GET · DELETE", "/admin/upload · /admin/runbooks"],
                ["GET",  "/admin/runbook-health · knowledge-gaps · feedback-stats"],
                ["GET · POST", "/auth/me · /auth/verify · /health"],
            ], size=10, header_size=10)
footer(s)


# ── Slide 10: Results · Roadmap · Close ────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG)
ambient(s, Inches(-2), Inches(4), Inches(7), color=AMBER, opacity_pct=7)
header(s, 10, TOTAL, "Outcome", "Results, Roadmap & Conclusion")

# Three columns
col_w = Inches(4.0); col_y = Inches(1.45); col_h = Inches(5.0)
gap_c = Inches(0.18)

# Benchmarks
card(s, Inches(0.5), col_y, col_w, col_h)
add_text(s, Inches(0.7), col_y + Inches(0.18), Inches(3.7), Inches(0.4),
         "Benchmarks", size=15, bold=True, color=AMBER)
build_table(s, Inches(0.7), col_y + Inches(0.7), Inches(3.65), Inches(3.8),
            ["Metric", "Result"], [
                ["Fast — TTFT",        "< 1 s"],
                ["Standard — E2E",     "3 – 5 s"],
                ["Deep — E2E (k=20)",  "5 – 9 s"],
                ["Cache-hit",          "< 50 ms"],
                ["FAISS top-12",       "< 10 ms"],
                ["Stream throughput",  "50 – 80 tok/s"],
            ], size=10, header_size=10)

# Roadmap
card(s, Inches(0.5) + col_w + gap_c, col_y, col_w, col_h)
add_text(s, Inches(0.7) + col_w + gap_c, col_y + Inches(0.18),
         Inches(3.7), Inches(0.4),
         "Future Roadmap", size=15, bold=True, color=AMBER)
add_bullets(s, Inches(0.7) + col_w + gap_c,
            col_y + Inches(0.7), Inches(3.65), Inches(4.2),
            [
                "Workspaces + RBAC (replace email whitelist).",
                "Compare / Diff mode (admin vs mine, A vs B).",
                "URL · Confluence · Notion · GitHub sync.",
                "OCR for scanned PDFs.",
                "Citation faithfulness verifier (2nd pass).",
                "Light/system theme · Settings page · PWA.",
                "Eval harness (recall@5 · MRR · grounding).",
            ], size=11, color=TEXT, accent=AMBER)

# Conclusion
card(s, Inches(0.5) + 2*(col_w + gap_c), col_y, col_w, col_h)
add_text(s, Inches(0.7) + 2*(col_w + gap_c), col_y + Inches(0.18),
         Inches(3.7), Inches(0.4),
         "Conclusion", size=15, bold=True, color=AMBER)
add_bullets(s, Inches(0.7) + 2*(col_w + gap_c),
            col_y + Inches(0.7), Inches(3.65), Inches(4.2),
            [
                "Citation-grounded RAG with strict gate-check.",
                "Six modes from one pipeline (~120 LoC).",
                "Hybrid retrieval: semantic + lexical + cross-encoder.",
                "Streaming UX with abort + threading.",
                "First-class artifacts: bookmark · share · export.",
                "Admin loop: knowledge gaps + runbook health.",
                "Modular: adding a mode is a dict entry.",
            ], size=11, color=TEXT, accent=EMERALD)

# Closing strap
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         "Thank you  ·  Questions?",
         size=24, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
footer(s)


# ── Save ───────────────────────────────────────────────────────────
out_path = r"C:/Users/ABHINAV TEJA/Downloads/Gen ai/resolveit-ai/ResolveIT_AI_Presentation_Themed.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
